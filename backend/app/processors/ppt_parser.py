"""
PowerPoint Parser.

Extracts:
- Text
- Tables
- Embedded Images

Each slide becomes a logical page.
"""

import io
from pathlib import Path
from typing import List, Tuple

from PIL import Image
from pptx import Presentation

from backend.app.processors.base import ParsedPage, RawDocumentElement
from backend.app.core.logging import logger


class PPTParser:

    def parse(
        self,
        file_path: Path
    ) -> Tuple[List[ParsedPage], List[RawDocumentElement], dict]:

        if not file_path.exists():
            raise FileNotFoundError(f"PPTX not found: {file_path}")

        prs = Presentation(str(file_path))

        pages: List[ParsedPage] = []
        elements: List[RawDocumentElement] = []

        for slide_idx, slide in enumerate(prs.slides):

            page_num = slide_idx + 1

            pages.append(
                ParsedPage(
                    page_number=page_num,
                    width=float(prs.slide_width),
                    height=float(prs.slide_height),
                    image=None,
                    raw_text=None
                )
            )

            reading_order = 1

            for shape in slide.shapes:

                bbox = [
                    float(shape.left),
                    float(shape.top),
                    float(shape.left + shape.width),
                    float(shape.top + shape.height)
                ]

                # -----------------------
                # TEXT
                # -----------------------
                if hasattr(shape, "text") and shape.text.strip():

                    elements.append(
                        RawDocumentElement(
                            type="text",
                            page=page_num,
                            bbox=bbox,
                            text=shape.text.strip(),
                            confidence=1.0,
                            attributes={
                                "reading_order": reading_order,
                                "source": "ppt_text"
                            }
                        )
                    )

                    reading_order += 1

                # -----------------------
                # TABLE
                # -----------------------
                if getattr(shape, "has_table", False):

                    table = shape.table

                    rows_data = []

                    for row in table.rows:
                        rows_data.append(
                            [cell.text.strip() for cell in row.cells]
                        )

                    if rows_data:

                        headers = rows_data[0]

                        md = "| " + " | ".join(headers) + " |\n"
                        md += "| " + " | ".join(["---"] * len(headers)) + " |\n"

                        for row in rows_data[1:]:
                            md += "| " + " | ".join(row) + " |\n"

                        html = (
                            "<table><thead><tr>"
                            + "".join(f"<th>{h}</th>" for h in headers)
                            + "</tr></thead><tbody>"
                        )

                        for row in rows_data[1:]:
                            html += (
                                "<tr>"
                                + "".join(f"<td>{c}</td>" for c in row)
                                + "</tr>"
                            )

                        html += "</tbody></table>"

                        elements.append(
                            RawDocumentElement(
                                type="table",
                                page=page_num,
                                bbox=bbox,
                                text="\n".join(
                                    [" | ".join(r) for r in rows_data]
                                ),
                                markdown=md,
                                html=html,
                                table_data={"rows": rows_data},
                                confidence=1.0,
                                attributes={
                                    "reading_order": reading_order,
                                    "source": "ppt_table"
                                }
                            )
                        )

                        reading_order += 1

                # -----------------------
                # IMAGE
                # -----------------------
                if shape.shape_type == 13:

                    try:
                        img_blob = shape.image.blob

                        pil_img = Image.open(
                            io.BytesIO(img_blob)
                        ).convert("RGB")

                        elements.append(
                            RawDocumentElement(
                                type="image",
                                page=page_num,
                                bbox=bbox,
                                image=pil_img,
                                confidence=1.0,
                                attributes={
                                    "reading_order": reading_order,
                                    "source": "ppt_image"
                                }
                            )
                        )

                        reading_order += 1

                    except Exception as e:
                        logger.warning(
                            f"Could not extract PPT image: {e}"
                        )

        metadata = {
            "title": file_path.stem,
            "page_count": len(prs.slides),
            "slides_count": len(prs.slides)
        }

        logger.info(
            f"Parsed PPTX {file_path.name}: "
            f"{len(prs.slides)} slides, "
            f"{len(elements)} elements"
        )

        return pages, elements, metadata


ppt_parser = PPTParser()